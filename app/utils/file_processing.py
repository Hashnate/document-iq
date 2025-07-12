import os
import magic
from typing import Optional
import pdfminer.high_level
from pptx import Presentation
from bs4 import BeautifulSoup
import docx
import csv
import openpyxl

def extract_text_from_file(file_path: str) -> Optional[str]:
    """Extract text from various file types"""
    mime = magic.Magic(mime=True)
    file_type = mime.from_file(file_path)
    
    try:
        if file_type == 'text/plain':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        
        elif file_type == 'application/pdf':
            return pdfminer.high_level.extract_text(file_path)
        
        elif file_type in ['application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                          'application/msword']:
            doc = docx.Document(file_path)
            return '\n'.join([para.text for para in doc.paragraphs])
        
        elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
        
        elif file_type == 'text/html':
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                return soup.get_text()
        
        elif file_type == 'text/csv':
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                return '\n'.join([','.join(row) for row in reader])
        
        elif file_type in ['application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                          'application/vnd.ms-excel']:
            wb = openpyxl.load_workbook(file_path)
            text = []
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    text.append('\t'.join([str(cell) if cell is not None else '' for cell in row]))
            return '\n'.join(text)
        
        else:
            print(f"Unsupported file type: {file_type}")
            return None
    
    except Exception as e:
        print(f"Error processing {file_path}: {str(e)}")
        return None